import streamlit as st
import seaborn as sns
import pandas as pd
import PIL.Image
import os
import tempfile
import google.generativeai as genai
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader
import json
import plotly.express as px
from docx import Document
from pptx import Presentation

# hàm lấy text từ file excel
def read_docx(file_path):
    doc = Document(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text
# hàm lấy text từ file powper point
def read_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

api_key = st.secrets['gemini']['GOOGLE_API_KEY']

if not api_key:
    st.error('Không tìm thấy key')
    st.stop()

genai.configure(api_key=api_key)

# helper functions Report Generator

def plot_chart(folder_path, chart_type, data, x_col, y_col):
    if chart_type == "Line Chart":
        chart = sns.lineplot(data=data, x=x_col, y=y_col, markers='o')
        # st.pyplot()
    elif chart_type == "Bar Chart":
        chart = sns.barplot(data=data, x=x_col, y=y_col)
        # st.pyplot()
    elif chart_type == "Scatter Plot":
        chart = sns.scatterplot(data=data, x=x_col, y=y_col)
        # st.pyplot()
    elif chart_type == "Pie Chart":
        pie_data = data.groupby(x_col)[y_col].sum()
        chart = pie_data.plot.pie(autopct='%1.1f%%', startangle=90)
        # st.pyplot()
    chart_fig = chart.get_figure()
    chart_name = f"{chart_type}_{x_col}_by_{y_col}.png"
    chart_fig.savefig(f"./{folder_path}/{chart_name}")

    chart_path = f"C:/Users/HP/Desktop/CS Data 08/final_project/charts/{chart_name}"

    return chart_path, chart_name

def plot_chart_interactions(folder_path, data, x_col, y_col):
    # Tạo biểu đồ
    chart_fig = px.scatter(data, x=x_col, y=y_col, title=f'Interaction between features {x_col} và {y_col}')
    
    # Đảm bảo thư mục tồn tại
    os.makedirs(folder_path, exist_ok=True)

    # Tạo đường dẫn file
    chart_name = f"Interactions_{x_col}_by_{y_col}.png"
    chart_path = f"C:/Users/HP/Desktop/CS Data 08/final_project/charts/{chart_name}"

    # Ghi ảnh
    chart_fig.write_image(chart_path)

    return chart_path, chart_name



def remove_chart(file_path):
    if os.path.isfile(file_path):
        os.remove(file_path)


def generate_report_from_chart(chart_folder, chart_name):
    genai.configure(api_key=st.secrets['gemini']['GEMINI_API_KEY']) # Thay key tại đây
    model = genai.GenerativeModel("gemini-2.0-flash")

    if chart_name.endswith(('.png', '.jpg', '.jpeg')):
        file_path = os.path.join(chart_folder, chart_name)
        organ = PIL.Image.open(file_path)

        response = model.generate_content(["According to the chart, generate a comprehensive report. Thus, what insight could be taken from it. Write at most 100 words", organ])
        bot_response = response.text.replace("*", "")

    return bot_response


def generate_excel_report(data, reports, report_name, save_option):
    """
    Generates an Excel report containing:
    - Original dataset
    - Pivot tables from aggregation
    - Charts (saved as images)
    - AI-generated insights

    Args:
    - data (pd.DataFrame): The original dataset.
    - reports (list): A list of reports with pivot tables, charts, and insights.
    - report_name (str): The output Excel file name (without extension).

    Returns:
    - str: The absolute path to the generated report.
    """

    report_filename = f"{report_name}.xlsx"
    report_path = os.path.abspath(report_filename)

    # Use Pandas' ExcelWriter with XlsxWriter as the engine
    with pd.ExcelWriter(report_path, engine='xlsxwriter') as writer:
        workbook = writer.book

        # Save the original dataset in the first sheet
        data.to_excel(writer, sheet_name="Datasource", index=False)
        # Lưu data overview
        if save_option['statistics']:
            overview_dict = {
                            'Number of variables': data.shape[1],
                            'Number of observations': data.shape[0],
                            'Missing cells': data.isnull().sum().sum(),
                            'Missing cells (%)': round(data.isnull().sum().sum() / (data.shape[0]*data.shape[1]) * 100, 2) if data.shape[0]*data.shape[1] > 0 else 0,
                            'Duplicate rows': data.duplicated().sum(),
                            'Duplicate rows (%)': round(data.duplicated().sum() / data.shape[0] * 100, 2) if data.shape[0] > 0 else 0,
                        }
            data_overview = pd.DataFrame(list(overview_dict.items()), columns=['Metric', 'Value'])
            data_overview['Value'] = data_overview['Value'].apply(lambda x: int(x) if isinstance(x, (int, float)) and x == int(x) else x)
            data_overview.to_excel(writer, sheet_name="Data_overview", index=False)
        if len(save_option['variables']):
            all_df = pd.DataFrame()
            for choose_feature in save_option['variables']:
                overview_dict = {
                    'Name': choose_feature,
                    'Distinct': data[choose_feature].nunique(),
                    'Missing cells': data[choose_feature].isnull().sum().sum(),
                    'Missing cells (%)': round(data[choose_feature].isnull().sum() / data.shape[1] * 100, 2),
                }
                data_overview = pd.DataFrame(list(overview_dict.items()), columns=['Metric', 'Value'])
                data_overview['Value'] = data_overview['Value'].apply(lambda x: int(x) if isinstance(x, (int, float)) and x == int(x) else x)
                all_df = pd.concat([all_df, data_overview], ignore_index=True)
            all_df.to_excel(writer, sheet_name="Variables", index=False)
        # Lưu data Correlations
        if save_option['statistics']:
            corr = data.corr(numeric_only=True)
            pd.DataFrame(corr).to_excel(writer, sheet_name="Correlations", index=False)
        # Process each report
        for report in reports:
            sheet_name = report["sheet_name"]

            # Save the pivot table to Excel
            report["pivot_table"].to_excel(writer, sheet_name=sheet_name, index=False)

            # Get the worksheet to add images and insights
            worksheet = writer.sheets[sheet_name]

            # Add Chart Image if Exists
            chart_path = report["chart_path"]
            if os.path.exists(chart_path):
                worksheet.insert_image("F1", chart_path, {"x_scale": 1, "y_scale": 1})
            cell_format = workbook.add_format({
                'align': 'center',     # Căn giữa ngang
                'valign': 'vcenter',   # Căn giữa dọc
                'text_wrap': True,     # Tự động xuống dòng
                'border': 1            # Viền xung quanh
            })
            # Add AI-Generated Insight in Cell F24
            worksheet.merge_range("F27:P33", report["insight"], cell_format)


# helper functions RAG
def get_file_text(uploaded_files, verbose=False):
    """
    Extract text from a list of uploaded files using appropriate LangChain loaders.
    
    Args:
        uploaded_files: List of file objects with name and read() method.
        verbose: If True, print detailed processing info.
    
    Returns:
        str: Concatenated text from all files.
    """
    all_text = ''
    for file in uploaded_files:
        try:
            suffix = os.path.splitext(file.name)[1].lower()
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp_file:
                tmp_file.write(file.read())
                tmp_file_path = tmp_file.name

            try:
                # Select loader based on file extension
                if suffix == '.pdf':
                    st.warning("tải file")
                    loader = PyPDFLoader(tmp_file_path)
                    # Extract text from the file
                    for page in loader.load_and_split():
                        all_text += page.page_content + '\n'
                    st.warning("tải file done")
                elif suffix in ['.doc', '.docx']:
                    st.warning("tải file")
                    all_text += read_docx(tmp_file_path)
                elif suffix in ['.ppt', '.pptx']:
                    all_text = read_pptx(tmp_file_path)
                else:
                    st.warning(f"⚠️ Unsupported file format: {suffix} for file {file.name}")
                    continue

                if verbose:
                    st.info(f"Processed file: {file.name}")

            finally:
                # Ensure temporary file is deleted
                try:
                    os.unlink(tmp_file_path)
                except Exception as e:
                    st.warning(f"Failed to delete temporary file {tmp_file_path}: {str(e)}")

        except Exception as e:
            st.error(f"Error processing file {file.name}: {str(e)}")
            continue
    st.warning("trả về tải file done")
    return all_text
@st.cache_data  
def get_text_chunk(text):
    try:
        # Giảm chunk_size và chunk_overlap để giảm tải cho API nhúng
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
        chunks = text_splitter.split_text(text)
        return chunks
    except Exception as e:
        # st.error là hàm của Streamlit, bạn có thể thay thế bằng print hoặc log tùy vào môi trường
        print(f'Lỗi chia chunk: {str(e)}')
        return []
    
@st.cache_resource
def get_vector_store(text_chunks):
    try:
        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
        
        # Chia các chunk thành các lô nhỏ (ví dụ: 100 chunk mỗi lô)
        batch_size = 100
        total_chunks = len(text_chunks)
        
        # Tạo vector store từ lô đầu tiên
        vector_store = FAISS.from_texts(text_chunks[:batch_size], embedding=embeddings)
        
        # Thêm các lô tiếp theo
        for i in range(batch_size, total_chunks, batch_size):
            batch_chunks = text_chunks[i:i+batch_size]
            vector_store.add_texts(batch_chunks)
            print(f"Đã xử lý xong lô từ {i} đến {i+batch_size}")
            
        vector_store.save_local("faiss_index")
    except Exception as e:
        print(f'Lỗi lưu vector database: {str(e)}')

def get_conversational_chain():
    prompt_template = """
    Trả lời câu hỏi một cách chi tiết nhất có thể dựa trên ngữ cảnh được cung cấp. Nếu câu trả lời không có trong ngữ cảnh được cung cấp, hãy nói, "Câu trả lời không có trong ngữ cảnh."
    Không cung cấp thông tin sai lệch.
    {history_block}
    Ngữ cảnh: {context}
    Câu hỏi: {question}

    Answer:
    """
    try:
        #temperature là độ sáng tạo, để 0.3 để focus vào tài liệu để ko lấy thông tin bên ngoài
        # nếu cho chỉ số cao thì mình phải finetune nó để nó làm theo ý mình
        model = ChatGoogleGenerativeAI(model='gemini-2.0-flash', temperature=0.3)
         # Thay {history_block} bằng chuỗi đã dựng
         #tạo conversation history từ session state
        chat_history = [
            {"role": msg["role"], "parts": [msg["content"] or '']}
            for msg in st.session_state.chat_history
        ]
        if chat_history:
            history_text = "Lịch sử hội thoại:\n"
            for item in chat_history:
                history_text += f"{item['role']}: {item['parts']}\n"
        else:
            history_text = ""
        
        prompt_template = prompt_template.replace("{history_block}", history_text)
        #context là tài liệu mình đưa vào
        #question là câu hỏi của user
        prompt = PromptTemplate(template=prompt_template, input_variables=['context','question'])
        #chain_type='stuff' lấy tất các chunks và question đưa vào mô hình
        chain = load_qa_chain(model, chain_type='stuff', prompt=prompt)
        return chain
    except Exception as e:
        st.error(f'Lỗi trong quá trình phân tích: {str(e)}')

def user_input(user_question):
    try:
        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
        if not os.path.exists('faiss_index'):
            st.error('Không tìm thấy FAISS index. Hãy tải file pdf lên trước')
            return
        #allow_dangerous_deserialization=False báo lỗi khi có mã độc, True thì bỏ qua
        new_db = FAISS.load_local('faiss_index', embeddings, allow_dangerous_deserialization=True)
        docs = new_db.similarity_search(user_question)
        st.warning("runing")
        chain = get_conversational_chain()
        st.warning('done')
        if not chain:
            return
        
        response = chain(
            {"input_documents": docs, "question": user_question},
            return_only_outputs=True
        )
        st.warning('done response')
        return response["output_text"]
    except Exception as e:
        st.error(f'Lỗi xử lý câu hỏi: {str(e)}')
        
# lưu đoạn chat
def load_history_chat():
    try:
        with open('chat_history.json', 'r', encoding='utf-8') as file:
            return json.load(file)
    except FileNotFoundError:
        return []

def save_chat_history():
    with open('chat_history.json', 'w', encoding='utf-8') as file:
        json.dump(st.session_state.chat_history, file, ensure_ascii=False, indent=2)